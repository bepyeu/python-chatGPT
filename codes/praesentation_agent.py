import streamlit as st
import pandas as pd
from pptx import Presentation
import tempfile
import os
from openai import OpenAI

class SimplePresentationAgent:
    def __init__(self):
        # LM Studio Connection
        self.client = OpenAI(base_url="http://localhost:1234/v1", api_key="not-needed")
    
    def create_simple_presentation(self, topic, slides_count=5):
        """Erstellt eine einfache Präsentation"""
        prs = Presentation()
        
        # Titelfolie
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = f"Präsentation: {topic}"
        title_slide.placeholders[1].text = "Erstellt mit AI Agent"
        
        # Inhaltsfolien
        for i in range(slides_count):
            content = self.generate_slide_content(topic, i+1)
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"{topic} - Teil {i+1}"
            slide.placeholders[1].text = content
        
        return prs
    
    def generate_slide_content(self, topic, slide_num):
        """Generiert Folieninhalt"""
        prompt = f"""
        Erstelle Inhalt für Präsentationsfolie {slide_num} über: {topic}
        
        Format:
        - 3-5 Bullet Points
        - Jeder Punkt prägnant
        - Professioneller Ton
        
        Gib NUR die Bullet Points zurück.
        """
        
        try:
            response = self.client.chat.completions.create(
                model="local-model",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=300,
                temperature=0.7
            )
            return response.choices[0].message.content
        except:
            return f"• Punkt 1 zu {topic}\n• Punkt 2 zu {topic}\n• Punkt 3 zu {topic}"

def main():
    st.set_page_config(
        page_title="Einfacher Präsentations Agent",
        page_icon="🤖",
        layout="centered"
    )
    
    st.title("🤖 Einfacher Präsentations-Agent")
    st.markdown("Erstelle schnell Präsentationen mit KI")
    
    # Agent initialisieren
    if 'agent' not in st.session_state:
        st.session_state.agent = SimplePresentationAgent()
    
    # Eingaben
    col1, col2 = st.columns(2)
    with col1:
        topic = st.text_input("Thema der Präsentation", "Künstliche Intelligenz")
    with col2:
        slides_count = st.slider("Anzahl Folien", 3, 10, 5)
    
    # Buttons
    col1, col2, col3 = st.columns(3)
    
    with col1:
        if st.button("🚀 Präsentation erstellen", type="primary"):
            with st.spinner("Erstelle Präsentation..."):
                prs = st.session_state.agent.create_simple_presentation(topic, slides_count)
                filename = f"{topic.replace(' ', '_')}_praesentation.pptx"
                prs.save(filename)
                
                with open(filename, "rb") as file:
                    st.download_button(
                        label="📥 Präsentation herunterladen",
                        data=file,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                st.success(f"✅ Präsentation mit {slides_count} Folien erstellt!")
    
    with col2:
        if st.button("🧹 Zurücksetzen"):
            st.session_state.clear()
            st.rerun()
    
    with col3:
        if st.button("ℹ️ Hilfe"):
            st.info("""
            **So funktioniert's:**
            1. Thema eingeben
            2. Folienanzahl wählen  
            3. Auf 'Präsentation erstellen' klicken
            4. Herunterladen und öffnen!
            """)
    
    # Chat-Funktion
    st.markdown("---")
    st.subheader("💬 Schnell-Chat")
    
    chat_input = st.text_input("Frage den Agent...")
    if chat_input:
        with st.spinner("Denke nach..."):
            response = st.session_state.agent.generate_slide_content(chat_input, 1)
            st.write(f"**Antwort:** {response}")

if __name__ == "__main__":
    main()